# Requirements: pip install easyocr pdf2image pillow tqdm docx2txt openpyxl extract_msg textract python-pptx beautifulsoup4




import logging
import subprocess
from pathlib import Path
from tempfile import TemporaryDirectory
from tqdm import tqdm
import docx2txt
import easyocr
import zipfile
from bs4 import BeautifulSoup
import openpyxl
import extract_msg
import email
from pdf2image import convert_from_path
from PIL import Image



# -----------------------------
# Paths
# -----------------------------
INPUT_DIR = Path(r"E:\Ermine\mywritingpad@proton.me\attachments (output)")
OUTPUT_DIR = Path(r"E:\Ermine\mywritingpad@proton.me\Testing\OCR-test")
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
LOG_FILE = OUTPUT_DIR / "conversion.log"


# -----------------------------
# Logging setup
# -----------------------------
logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s [%(levelname)s] %(message)s",
    handlers=[
        logging.FileHandler(LOG_FILE, encoding="utf-8"),
        logging.StreamHandler()
    ]
)


# -----------------------------
# OCR reader
# -----------------------------
reader = easyocr.Reader(['en', 'nl', 'es'])  # Add 'nl' or others if needed

def convert_pdf_to_text(file_path: Path, out_txt: Path):
    """Convert PDF to text using EasyOCR (no Tesseract)."""
    with TemporaryDirectory() as tempdir:
        tempdir = Path(tempdir)
        pdf_pages = convert_from_path(file_path, dpi=300)
        all_text = []
        for i, page in enumerate(pdf_pages, start=1):
            img_path = tempdir / f"page_{i:03}.jpg"
            page.save(img_path, "JPEG")
            try:
                text_blocks = reader.readtext(str(img_path), detail=0, paragraph=True)
                all_text.append("\n".join(text_blocks))
            except Exception as e:
                logging.warning(f"⚠️ Error OCRing page {i} of {file_path.name}: {e}")
        out_txt.write_text("\n\n".join(all_text), encoding="utf-8")

# -----------------------------
# Office / text converters
# -----------------------------
def convert_docx_to_text(file_path, out_txt):
    text = docx2txt.process(str(file_path))
    out_txt.write_text(text, encoding="utf-8")

def convert_image_to_text(file_path, out_txt):
    result = reader.readtext(str(file_path), detail=0, paragraph=True)
    text = "\n".join(result)
    out_txt.write_text(text, encoding="utf-8")

def convert_pptx_to_text(file_path, out_txt):
    from pptx import Presentation
    prs = Presentation(file_path)
    text = "\n".join(shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text"))
    out_txt.write_text(text, encoding="utf-8")

def convert_excel_to_text(file_path, out_txt):
    wb = openpyxl.load_workbook(file_path, data_only=True)
    text = ""
    for ws in wb.worksheets:
        for row in ws.iter_rows(values_only=True):
            text += " ".join(str(cell or "") for cell in row) + "\n"
    out_txt.write_text(text, encoding="utf-8")

def convert_eml_to_text(file_path, out_txt):
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        msg = email.message_from_file(f)
    parts = [msg.get("From", ""), msg.get("To", ""), msg.get("Subject", "")]
    for part in msg.walk():
        if part.get_content_type() == "text/plain":
            try:
                parts.append(part.get_payload(decode=True).decode("utf-8", errors="ignore"))
            except Exception:
                pass
    out_txt.write_text("\n".join(parts), encoding="utf-8")

def convert_msg_to_text(file_path, out_txt):
    msg = extract_msg.Message(str(file_path))
    text = f"From: {msg.sender}\nTo: {msg.to}\nSubject: {msg.subject}\n\n{msg.body}"
    out_txt.write_text(text, encoding="utf-8")

def convert_html_to_text(file_path, out_txt):
    soup = BeautifulSoup(file_path.read_text(encoding="utf-8", errors="ignore"), "html.parser")
    text = soup.get_text(separator="\n")
    out_txt.write_text(text, encoding="utf-8")

def convert_zip_to_text(file_path, out_txt):
    with zipfile.ZipFile(file_path, "r") as z:
        extract_dir = OUTPUT_DIR / "unzipped" / file_path.stem
        extract_dir.mkdir(parents=True, exist_ok=True)
        z.extractall(extract_dir)
    out_txt.write_text(f"Extracted {file_path.name} to {extract_dir}", encoding="utf-8")

#def convert_with_textract(file_path, out_txt):
#    text = textract.process(str(file_path)).decode("utf-8", errors="ignore")
#    out_txt.write_text(text, encoding="utf-8")

# -----------------------------
# Dispatcher
# -----------------------------
def process_file(file_path):
    ext = file_path.suffix.lower()
    out_txt = OUTPUT_DIR / f"{file_path.stem}.txt"
    logging.info(f"Processing: {file_path}")

    try:
        if ext == ".pdf":
            convert_pdf_to_text(file_path, out_txt)
        elif ext in [".docx", ".doc"]:
            convert_docx_to_text(file_path, out_txt)
        elif ext in [".jpg", ".jpeg", ".png", ".tif", ".gif"]:
            convert_image_to_text(file_path, out_txt)
        elif ext == ".pptx":
            convert_pptx_to_text(file_path, out_txt)
        elif ext in [".xlsx", ".xls", ".xlsm", ".xlsb"]:
            convert_excel_to_text(file_path, out_txt)
        elif ext == ".eml":
            convert_eml_to_text(file_path, out_txt)
        elif ext == ".msg":
            convert_msg_to_text(file_path, out_txt)
        elif ext in [".html", ".htm", ".xml", ".rtf", ".txt", ".csv"]:
            convert_html_to_text(file_path, out_txt)
        elif ext == ".zip":
            convert_zip_to_text(file_path, out_txt)
        else:
            convert_with_textract(file_path, out_txt)

        logging.info(f"✅ Success: {file_path.name}")

    except Exception as e:
        logging.error(f"❌ Error processing {file_path}: {e}")

# -----------------------------
# Main
# -----------------------------
def main():
    files = [f for f in INPUT_DIR.glob("**/*") if f.is_file()]
    logging.info(f"Found {len(files)} files in {INPUT_DIR}")
    for f in tqdm(files, desc="Processing files"):
        process_file(f)
    logging.info("All files processed.")

if __name__ == "__main__":
    main()
